"""T13 多格式 e2e（version_test 风格，直接调模块、需干净库）。
逐格式：DOCX/PPTX/XLSX/HTML（+PDF，需 reportlab）→ ingest → chunks>0 → 检索命中 marker → position 行为。
运行：.venv/bin/python scripts/multiformat_test.py"""
import hashlib
import io
import os
import sys
import uuid

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from app.bootstrap import default_tenant_id, default_user_id
from app.config import settings
from app.db import get_conn
from app.ingest import pipeline
from app.storage import get_minio

TID = default_tenant_id()
UID = default_user_id()


def _docx(marker):
    import docx

    d = docx.Document()
    d.add_heading("Report Section", level=1)
    d.add_paragraph(" ".join([marker] * 30))
    bio = io.BytesIO()
    d.save(bio)
    return bio.getvalue()


def _pptx(marker):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deck Title"
    slide.placeholders[1].text = " ".join([marker] * 30)
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _xlsx(marker):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["col"])
    for _ in range(30):
        ws.append([marker])
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def _html(marker):
    return ("<h1>Page Title</h1><p>" + (marker + " ") * 30 + "</p>").encode()


def _pdf(marker):
    from reportlab.pdfgen import canvas

    bio = io.BytesIO()
    c = canvas.Canvas(bio)
    c.drawString(50, 700, " ".join([marker] * 5))
    c.showPage()
    c.save()
    return bio.getvalue()


FORMATS = [
    ("DOCX", _docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "x.docx"),
    ("PPTX", _pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation", "x.pptx"),
    ("XLSX", _xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "x.xlsx"),
    ("HTML", _html, "text/html", "x.html"),
]


def _setup(file_id, kb_id, data, mime, name):
    sk = f"{file_id}/v1/raw"
    get_minio().put_object(settings.minio_bucket, sk, io.BytesIO(data), len(data))
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("INSERT INTO kb_kb(id,tenant_id,name,owner_id) VALUES(%s,%s,%s,%s)", (kb_id, TID, f"mf-{kb_id[:8]}", UID))
            cur.execute(
                "INSERT INTO kb_file(id,tenant_id,storage_key,name,content_hash,mime,status,owner_user_id) "
                "VALUES(%s,%s,%s,%s,%s,%s,'parsing',%s)",
                (file_id, TID, sk, name, hashlib.sha256(data).hexdigest(), mime, UID),
            )
            cur.execute("INSERT INTO kb_file_kb(file_id,kb_id,tenant_id) VALUES(%s,%s,%s)", (file_id, kb_id, TID))


def _search_hits(query):
    from app.middleware.auth import Principal
    from app.retrieval.orchestrator import retrieve

    return " ".join(h["snippet"] for h in retrieve(query, Principal(TID, UID))["hits"])


def _position_null(file_id):
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT count(*) FROM kb_chunk WHERE file_id=%s AND position IS NOT NULL", (file_id,))
            return cur.fetchone()[0] == 0


def _position_not_null(file_id):
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT count(*) FROM kb_chunk WHERE file_id=%s AND position IS NOT NULL", (file_id,))
            return cur.fetchone()[0] > 0


def main():
    fails = []

    def check(name, cond, detail=""):
        print(f"  [{'PASS' if cond else 'FAIL'}] {name}" + (f" :: {detail}" if detail else ""))
        if not cond:
            fails.append(name)

    print("逐格式 ingest（DOCX/PPTX/XLSX/HTML）…")
    for label, fn, mime, name in FORMATS:
        marker = f"MK_{label}_" + uuid.uuid4().hex[:6]
        data = fn(marker)
        fid = str(uuid.uuid4())
        kb = str(uuid.uuid4())
        _setup(fid, kb, data, mime, name)
        st = pipeline.ingest_file(fid)
        check(f"{label} chunks>0", st["chunks"] > 0, f"chunks={st['chunks']}")
        check(f"{label} 检索命中 marker", marker in _search_hits(marker), "no hit")
        check(f"{label} position 全 NULL（Office 无 bbox）", _position_null(fid), "position unexpectedly set")

    print("PDF（需 reportlab）…")
    try:
        import reportlab  # noqa: F401
        marker = "MK_PDF_" + uuid.uuid4().hex[:6]
        data = _pdf(marker)
        fid = str(uuid.uuid4())
        kb = str(uuid.uuid4())
        _setup(fid, kb, data, "application/pdf", "x.pdf")
        st = pipeline.ingest_file(fid)
        check("PDF chunks>0", st["chunks"] > 0, f"chunks={st['chunks']}")
        check("PDF position 非 NULL（有 bbox）", _position_not_null(fid), "position null")
    except ImportError:
        check("PDF 跳过（无 reportlab）", True, "skipped")

    print(f"\n{'ALL GREEN ✅' if not fails else 'FAILURES ❌ ' + str(fails)}")
    sys.exit(1 if fails else 0)


if __name__ == "__main__":
    main()
