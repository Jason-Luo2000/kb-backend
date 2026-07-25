"""T13 parser 单测（纯函数，无需 DB/MinIO）。各格式用库内存构造 bytes，不提交二进制夹具。
运行：.venv/bin/pytest tests/test_parser.py -q"""
import io

import pytest

from app.adapters import parser
from app.adapters.parser import Block


# ---------- 夹具构造（内存 bytes）----------
def _docx_bytes():
    import docx

    d = docx.Document()
    d.add_heading("Doc Title", level=1)
    d.add_paragraph("Hello body text")
    tbl = d.add_table(rows=2, cols=2)
    tbl.rows[0].cells[0].text, tbl.rows[0].cells[1].text = "k1", "k2"
    tbl.rows[1].cells[0].text, tbl.rows[1].cells[1].text = "v1", "v2"
    bio = io.BytesIO()
    d.save(bio)
    return bio.getvalue()


def _pptx_bytes():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "bullet body text"
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _xlsx_bytes():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["a", 1])
    ws.append(["b", 2])
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


# ---------- MD / D0.2 ----------
def test_md_splits_title_body():
    b = parser.parse_bytes(b"# H1\nbody1\n\n## H2\nbody2", "text/markdown", "x.md")
    types = [x.block_type for x in b]
    assert types == ["title", "text", "title", "text"]
    assert b[0].text == "# H1" and b[0].level == 1 and b[0].section_path == "H1"
    assert b[2].text == "## H2" and b[2].level == 2 and b[2].section_path == "H2"


def test_md_chunk_content_byte_identity():
    """D0.2：naive_merge 重连 title+body → 与旧 per-section block 逐字一致。"""
    from app.ingest import chunker

    blocks = parser.parse_bytes(b"# H1\nbody1 line\n\n## H2\nbody2", "text/markdown", "x.md")
    chunks = chunker.chunk_blocks(blocks, size=512)
    contents = [c["content"] for c in chunks]
    assert contents == ["# H1\nbody1 line", "## H2\nbody2"]


# ---------- DOCX ----------
def test_docx_title_text_table():
    b = parser.parse_bytes(_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "x.docx")
    types = [x.block_type for x in b]
    assert "title" in types and "text" in types and "table" in types
    title = next(x for x in b if x.block_type == "title")
    assert title.text == "Doc Title" and title.level == 1


# ---------- PPTX ----------
def test_pptx_title_text():
    b = parser.parse_bytes(_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "x.pptx")
    types = [x.block_type for x in b]
    assert "title" in types and "text" in types
    assert all(x.page == 1 for x in b)
    title = next(x for x in b if x.block_type == "title")
    assert title.text == "Slide Title"


# ---------- XLSX ----------
def test_xlsx_table_blocks():
    b = parser.parse_bytes(_xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "x.xlsx")
    assert b, "xlsx should yield blocks"
    assert all(x.block_type == "table" for x in b)
    assert all(x.section_path == "Sheet1" for x in b)
    joined = "\n".join(x.text for x in b)
    assert "Name" in joined and "a" in joined


# ---------- HTML ----------
def test_html_title_text_table():
    html = b"<h1>T1</h1><p>hello world</p><table><tr><th>x</th></tr><tr><td>1</td></tr></table>"
    b = parser.parse_bytes(html, "text/html", "x.html")
    types = [x.block_type for x in b]
    assert "title" in types and "text" in types and "table" in types
    title = next(x for x in b if x.block_type == "title")
    assert title.text == "T1" and title.level == 1


# ---------- PDF（reportlab 生成真 PDF，无 reportlab 则跳过）----------
def test_pdf_text_block_has_bbox():
    pytest.importorskip("reportlab")
    import io
    from reportlab.pdfgen import canvas

    bio = io.BytesIO()
    c = canvas.Canvas(bio)
    c.drawString(50, 700, "Hello PDF marker text")
    c.showPage()
    c.save()
    b = parser.parse_bytes(bio.getvalue(), "application/pdf", "x.pdf")
    assert b, "pdf should yield a text block"
    txt = next(x for x in b if x.block_type == "text")
    assert txt.page == 1 and txt.bbox is not None  # T13：PDF 文本块带页级 bbox
    assert "marker" in txt.text


# ---------- 注册表分派 ----------
def test_registry_dispatch_by_ext():
    """mime 缺失时按扩展名分派。"""
    b = parser.parse_bytes(b"# md\nbody", None, "readme.md")
    assert any(x.block_type == "title" for x in b)


def test_unknown_falls_back_to_text():
    b = parser.parse_bytes(b"# h\nplain", "application/octet-stream", "x.unknown")
    assert b and b[0].block_type in ("title", "text")
