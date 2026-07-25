"""文档解析适配器（T13）：Parser Registry 按 mime/扩展名分派到各格式 handler。

输出统一「版式块序列」Block{page, text, section_path, block_type, bbox, level}：
- block_type ∈ title|text|table|figure|equation|header|footer|caption（驱动 naive_merge 边界 + skip_summary）
- bbox=(left,top,right,bottom) 仅 PDF 真能给（页级/表格）；Office 文本 bbox=None → position=NULL（page 仍填）
- level=标题深度 1..6

格式：PDF（pdfplumber + find_tables + OCR stub）、DOCX（python-docx）、PPTX（python-pptx）、
XLSX（openpyxl）、HTML（bs4+lxml）、MD/TXT（手写 # 分割）。OCR 默认关、lazy import（本地无 tesseract 时
warn+丢页，real OCR 走 Dockerfile 后续）。DeepDOC 版式/表格 TSR、FACTORY 多分块器 defer。
"""
import io
import os
import re
from dataclasses import dataclass

import pdfplumber

from app.config import settings


@dataclass
class Block:
    page: int
    text: str
    section_path: str | None = None
    block_type: str = "text"  # title|text|table|figure|equation|header|footer|caption
    bbox: tuple[float, float, float, float] | None = None  # (left, top, right, bottom)
    level: int | None = None  # 标题深度 1..6


def parse_bytes(data: bytes, mime: str | None, name: str) -> list[Block]:
    """按 mime → 扩展名 → 默认(MD/TXT) 分派。每个 handler 内部 lazy-import 自己的库。"""
    n = (name or "").lower()
    handler = _HANDLERS.get((mime or "").lower()) or _EXT.get(_ext(n))
    if handler is None:
        handler = _parse_text
    return handler(data)


def _ext(name: str) -> str:
    return os.path.splitext(name)[1]


# ============ PDF ============
def _parse_pdf(data: bytes) -> list[Block]:
    from app.adapters import ocr

    blocks: list[Block] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_bbox = (0.0, 0.0, float(page.width), float(page.height))
            # 表格（尽力富化，非 DeepDOC TSR）：find_tables
            try:
                tables = page.find_tables()
            except Exception:  # noqa: BLE001
                tables = []
            for t in tables:
                rows = t.extract() or []
                txt = "\n".join("\t".join("" if c is None else str(c) for c in r) for r in rows)
                if txt.strip():
                    blocks.append(
                        Block(page=i, text=txt, block_type="table",
                              bbox=tuple(t.bbox) if t.bbox else page_bbox)
                    )
            # 文本
            txt = (page.extract_text() or "").strip()
            if len(txt) >= settings.ocr_min_chars_per_page:
                blocks.append(Block(page=i, text=txt, block_type="text", bbox=page_bbox))
            elif settings.ocr_enabled:  # 扫描页 + OCR 开 → lazy OCR（本地无 tesseract 则 warn+丢页）
                ocr_txt = ocr.ocr_page(data, i).strip()
                if ocr_txt:
                    blocks.append(Block(page=i, text=ocr_txt, block_type="text", bbox=page_bbox))
            # else：扫描页且 OCR 关 → 丢页
    return blocks


# ============ DOCX ============
def _parse_docx(data: bytes) -> list[Block]:
    import docx

    doc = docx.Document(io.BytesIO(data))
    blocks: list[Block] = []
    stack: list[str] = []

    def push_heading(text: str, lvl: int | None):
        nonlocal stack
        if lvl:
            stack = stack[: lvl - 1]
        stack.append(text)

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = ((para.style.name if para.style else "") or "").lower()
        if style.startswith("heading") or style.startswith("title"):
            parts = style.split()
            lvl = int(parts[1]) if len(parts) > 1 and parts[1].isdigit() else (len(stack) + 1)
            push_heading(text, lvl)
            blocks.append(Block(page=1, text=text, section_path=" / ".join(stack),
                                block_type="title", level=lvl))
        else:
            blocks.append(Block(page=1, text=text, section_path=" / ".join(stack) or None))
    for tbl in doc.tables:
        rows = [["" if c is None else (c.text or "") for c in r.cells] for r in tbl.rows]
        blocks.extend(_table_blocks(rows, " / ".join(stack) or None, 1))
    return blocks


# ============ PPTX ============
def _parse_pptx(data: bytes) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = (shape.text_frame.text or "").strip()
                if not txt:
                    continue
                if title is None and shape == title_shape:
                    title = txt
                    blocks.append(Block(page=i, text=txt, section_path=txt, block_type="title"))
                else:
                    blocks.append(Block(page=i, text=txt, section_path=title))
            if shape.has_table:
                rows = [["" if c is None else (c.text or "") for c in r.cells] for r in shape.table.rows]
                blocks.extend(_table_blocks(rows, title, i))
    return blocks


# ============ XLSX ============
XLSX_MAX_CELLS = 200_000  # 防巨表爆炸


def _parse_xlsx(data: bytes) -> list[Block]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[Block] = []
    page = 0
    try:
        for ws in wb.worksheets:
            page += 1
            rows: list[list[str]] = []
            cells = 0
            for r in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v) for v in r]
                if any(v.strip() for v in vals):
                    rows.append(vals)
                    cells += len(vals)
                if cells >= XLSX_MAX_CELLS:
                    break
            if rows:
                blocks.extend(_table_blocks(rows, ws.title, page))
    finally:
        wb.close()
    return blocks


# ============ HTML ============
def _parse_html(data: bytes) -> list[Block]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "lxml")
    for t in soup(["script", "style"]):
        t.decompose()
    blocks: list[Block] = []
    stack: list[str] = []
    for el in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
        if el.name == "table":
            rows = []
            for tr in el.find_all("tr"):
                cells = [c.get_text(" ", strip=True) for c in tr.find_all(["td", "th"])]
                if any(cells):
                    rows.append(cells)
            blocks.extend(_table_blocks(rows, " / ".join(stack) or None, 1))
        elif el.name and el.name[0] == "h":
            lvl = int(el.name[1])
            txt = el.get_text(" ", strip=True)
            if not txt:
                continue
            stack = stack[: lvl - 1]
            stack.append(txt)
            blocks.append(Block(page=1, text=txt, section_path=" / ".join(stack),
                                block_type="title", level=lvl))
        else:  # p / li
            txt = el.get_text(" ", strip=True)
            if txt:
                blocks.append(Block(page=1, text=txt, section_path=" / ".join(stack) or None))
    return blocks


def _table_blocks(rows: list[list[str]], section_path: str | None, page: int) -> list[Block]:
    """表格行按 parser_table_rows_per_chunk 切成多个 table block（防巨表爆炸）。"""
    out: list[Block] = []
    n = max(1, settings.parser_table_rows_per_chunk)
    for i in range(0, len(rows), n):
        chunk = rows[i : i + n]
        txt = "\n".join("\t".join(c for c in r) for r in chunk)
        if txt.strip():
            out.append(Block(page=page, text=txt, section_path=section_path, block_type="table"))
    return out


# ============ MD / 纯文本 ============
def _parse_text(data: bytes) -> list[Block]:
    """Markdown/纯文本：每个 # 标题段发 [title(原始行), body] 两块。
    title.text 保留原始 '# ...' 行；naive_merge 用 \\n 重连 → 与旧实现逐字一致（保 content_hash/T12 复用，D0.2）。"""
    text = data.decode("utf-8", "ignore")
    blocks: list[Block] = []
    cur_section: str | None = None
    cur_level: int | None = None
    body: list[str] = []

    def flush_body():
        nonlocal body
        if body:
            joined = "\n".join(body).strip()
            if joined:
                blocks.append(Block(page=1, text=joined, section_path=cur_section))
            body = []

    for line in text.splitlines():
        s = line.strip()
        if s.startswith("#"):
            flush_body()
            hashes = len(s) - len(s.lstrip("#"))
            cur_section = s.lstrip("#").strip() or cur_section
            cur_level = hashes or None
            blocks.append(Block(page=1, text=s, section_path=cur_section, block_type="title", level=cur_level))
        else:
            body.append(line)
    flush_body()
    return blocks or [Block(page=1, text=text.strip())]


# ============ 注册表（定义于 handler 之后）============
_HANDLERS = {
    "application/pdf": _parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _parse_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _parse_pptx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _parse_xlsx,
    "text/html": _parse_html,
    "application/xhtml+xml": _parse_html,
    "text/markdown": _parse_text,
    "text/plain": _parse_text,
}
_EXT = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".md": _parse_text,
    ".markdown": _parse_text,
    ".txt": _parse_text,
}
